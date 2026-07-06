from __future__ import annotations

from datetime import date
from pathlib import Path
from textwrap import wrap

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


OUT = Path("/Users/WilliamChang/Documents/Playground/IterationPilot/artifacts/docs/NexaFlow_从提问到回答_技术处理链路_2026-07-03.pptx")
LATEST = Path("/Users/WilliamChang/Documents/Playground/IterationPilot/artifacts/docs/latest_ppt_workspace.txt")

LATIN_SERIF = "Tiempos Text"
CJK_SERIF = "Songti SC"
MONO = "Menlo"

PAPER = RGBColor(248, 245, 239)
INK = RGBColor(39, 35, 31)
MUTED = RGBColor(102, 94, 86)
FAINT = RGBColor(142, 132, 120)
RUST = RGBColor(172, 86, 48)
RUST_DARK = RGBColor(126, 55, 35)
SAGE = RGBColor(89, 112, 93)
BLUE = RGBColor(53, 87, 119)
BORDER = RGBColor(216, 207, 195)
CARD = RGBColor(255, 252, 247)
SOFT = RGBColor(241, 235, 226)
CODE_BG = RGBColor(31, 35, 40)
CODE_FG = RGBColor(241, 245, 249)


def set_run_font(run, size=16, bold=False, color=INK, italic=False, font=LATIN_SERIF, cjk_font=CJK_SERIF):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color

    rpr = run._r.get_or_add_rPr()
    for tag, typeface in (("a:latin", font), ("a:ea", cjk_font), ("a:cs", font)):
        node = rpr.find(qn(tag))
        if node is None:
            node = OxmlElement(tag)
            rpr.append(node)
        node.set("typeface", typeface)


def add_textbox(slide, x, y, w, h, text="", size=16, bold=False, color=INK, align=None, font=LATIN_SERIF, cjk_font=CJK_SERIF):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    if align is not None:
        p.alignment = align
    run = p.add_run()
    run.text = text
    set_run_font(run, size=size, bold=bold, color=color, font=font, cjk_font=cjk_font)
    return box


def add_title(slide, title, kicker="NexaFlow 技术链路", page=None):
    add_textbox(slide, 0.55, 0.35, 2.7, 0.28, kicker.upper(), size=8.5, bold=True, color=RUST_DARK)
    add_textbox(slide, 0.55, 0.66, 9.9, 0.62, title, size=27, bold=True, color=INK)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(1.34), Inches(12.15), Inches(0.018))
    line.fill.solid()
    line.fill.fore_color.rgb = BORDER
    line.line.fill.background()
    if page is not None:
        add_textbox(slide, 12.1, 7.12, 0.7, 0.18, f"{page:02d}", size=8, bold=True, color=FAINT, align=PP_ALIGN.RIGHT)


def add_footer(slide, page):
    add_textbox(slide, 0.55, 7.13, 5.4, 0.18, "NexaFlow｜Question → Answer Pipeline", size=7.2, color=FAINT)
    add_textbox(slide, 11.2, 7.13, 1.5, 0.18, f"{page:02d}", size=7.2, bold=True, color=FAINT, align=PP_ALIGN.RIGHT)


def add_card(slide, x, y, w, h, title, body=None, fill=CARD, accent=RUST, title_size=14.5, body_size=11.2):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = BORDER
    shape.line.width = Pt(0.7)
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(0.16)
    tf.margin_right = Inches(0.14)
    tf.margin_top = Inches(0.12)
    tf.margin_bottom = Inches(0.10)
    tf.word_wrap = True

    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    set_run_font(r, size=title_size, bold=True, color=accent)
    if body:
        for item in body:
            p = tf.add_paragraph()
            p.space_before = Pt(4)
            p.level = 0
            r = p.add_run()
            r.text = item
            set_run_font(r, size=body_size, color=INK)
    return shape


def add_bullets(slide, x, y, w, h, bullets, size=12.5, color=INK, leading=1.05, bullet=True):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.01)
    tf.margin_bottom = Inches(0.01)
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(5 * leading)
        p.level = 0
        r = p.add_run()
        r.text = ("• " if bullet else "") + item
        set_run_font(r, size=size, color=color)
    return box


def add_label(slide, x, y, w, h, text, fill=SOFT, color=RUST_DARK, size=10.5, bold=True):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = BORDER
    shape.line.width = Pt(0.45)
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    set_run_font(r, size=size, bold=bold, color=color)
    return shape


def add_code(slide, x, y, w, h, lines, title=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CODE_BG
    shape.line.color.rgb = RGBColor(69, 78, 90)
    shape.line.width = Pt(0.7)
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(0.14)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.10)
    tf.margin_bottom = Inches(0.10)
    if title:
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = title
        set_run_font(r, size=10.5, bold=True, color=RGBColor(250, 204, 121), font=MONO, cjk_font=MONO)
        first = False
    else:
        first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(1)
        r = p.add_run()
        r.text = line
        set_run_font(r, size=8.8, color=CODE_FG, font=MONO, cjk_font=MONO)
    return shape


def add_arrow(slide, x1, y1, x2, y2, color=RUST):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    line.line.color.rgb = color
    line.line.width = Pt(1.5)
    line.line.end_arrowhead = True
    return line


def blank(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = PAPER
    return slide


def add_table(slide, x, y, w, h, headers, rows, widths=None, font_size=8.8):
    rows_count = len(rows) + 1
    cols = len(headers)
    tbl = slide.shapes.add_table(rows_count, cols, Inches(x), Inches(y), Inches(w), Inches(h)).table
    if widths:
        for i, width in enumerate(widths):
            tbl.columns[i].width = Inches(width)
    for c, header in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = SOFT
        cell.text = header
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                set_run_font(r, size=font_size, bold=True, color=RUST_DARK)
    for r_idx, row in enumerate(rows, start=1):
        for c, text in enumerate(row):
            cell = tbl.cell(r_idx, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = CARD if r_idx % 2 else RGBColor(252, 248, 241)
            cell.text = text
            for p in cell.text_frame.paragraphs:
                p.space_after = Pt(0)
                for run in p.runs:
                    set_run_font(run, size=font_size, color=INK)
    return tbl


def cover(prs):
    slide = blank(prs)
    add_textbox(slide, 0.65, 0.52, 2.4, 0.28, "NEXAFLOW INTERNAL", size=9, bold=True, color=RUST_DARK)
    add_textbox(slide, 0.65, 1.35, 10.6, 1.25, "从提问到回答：\nAI 分析链路的技术处理细节", size=36, bold=True, color=INK)
    add_textbox(
        slide,
        0.72,
        3.05,
        8.4,
        0.8,
        "一份面向工程、产品和交付团队的端到端技术说明：用户输入如何变成可验证、可追溯、可渲染的业务回答。",
        size=17,
        color=MUTED,
    )
    add_label(slide, 0.72, 4.35, 2.0, 0.36, "Question Intake", fill=RGBColor(239, 229, 217))
    add_label(slide, 2.95, 4.35, 2.0, 0.36, "Context Build", fill=RGBColor(239, 229, 217))
    add_label(slide, 5.18, 4.35, 2.0, 0.36, "Harness", fill=RGBColor(239, 229, 217))
    add_label(slide, 7.41, 4.35, 2.0, 0.36, "Validation", fill=RGBColor(239, 229, 217))
    add_label(slide, 9.64, 4.35, 2.0, 0.36, "Answer UI", fill=RGBColor(239, 229, 217))
    for x in [2.74, 4.97, 7.20, 9.43]:
        add_arrow(slide, x, 4.53, x + 0.18, 4.53, color=RUST_DARK)
    add_textbox(slide, 0.72, 6.25, 5.0, 0.3, f"Generated {date.today().isoformat()} · Typeface: Tiempos Text / Songti SC fallback", size=8.5, color=FAINT)
    return slide


def build_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    cover(prs)

    page = 2

    slide = blank(prs)
    add_title(slide, "这条链路解决的不是“问大模型”这么简单", page=page)
    add_card(slide, 0.65, 1.72, 3.7, 4.8, "普通聊天壳", [
        "用户问题直接拼进 prompt。",
        "模型自己理解表格、口径和证据。",
        "输出后很难判断数字是否来自本地数据。",
        "追问时上下文容易漂移。"
    ], fill=RGBColor(255, 250, 244), accent=FAINT)
    add_card(slide, 4.85, 1.72, 7.35, 4.8, "NexaFlow 的处理方式", [
        "用户消息进入会话系统后，先形成可持久化 job，而不是同步阻塞 UI。",
        "系统构建数据覆盖、业务空间、表结构、知识库、外部证据和历史纠偏上下文。",
        "路由器判断：轻量追问、缓存追问、完整重算、报告生成、强 Harness。",
        "涉及数字和证据时，本地计算和 ValidationDecisionEngine 对回答做约束。",
        "最终回答附带 evidence、coverage、trace、job record，能复盘。"
    ], fill=CARD, accent=RUST)
    add_footer(slide, page); page += 1

    slide = blank(prs)
    add_title(slide, "端到端总览：从用户输入到屏幕上的回答", page=page)
    stages = [
        ("1 用户输入", "Composer\n引用 / 范围 / 资料策略"),
        ("2 会话落库", "User message\nPersistentAIJob"),
        ("3 上下文准备", "Coverage\nTables\nBusiness space"),
        ("4 路由决策", "Quick / Cached\nFull / Harness"),
        ("5 执行", "AI Streaming\nLocal metrics\nSQL/Notebook"),
        ("6 校验", "Trace numbers\nIssue classify\nRepair"),
        ("7 渲染", "Markdown\nEvidence links\nAuto-scroll")
    ]
    x = 0.45
    for i, (title, body) in enumerate(stages):
        add_card(slide, x, 1.85, 1.65, 2.4, title, body.split("\n"), fill=CARD, accent=[RUST, BLUE, SAGE][i % 3], title_size=11.5, body_size=8.6)
        if i < len(stages) - 1:
            add_arrow(slide, x + 1.65, 3.0, x + 1.93, 3.0, color=FAINT)
        x += 1.82
    add_textbox(slide, 0.72, 5.15, 11.7, 0.55, "关键设计：每一步都产生结构化中间物，后续阶段只消费明确契约，而不是依赖一整段不可控 prompt。", size=17, bold=True, color=INK)
    add_bullets(slide, 0.9, 5.85, 11.2, 0.72, [
        "UI 负责收集意图；Store 负责编排任务；Services 负责事实层和 AI 调用；Support/Views 负责最终展示和性能。"
    ], size=11.5)
    add_footer(slide, page); page += 1

    slide = blank(prs)
    add_title(slide, "1. 提问入口：Composer 收集的不只是文本", page=page)
    add_card(slide, 0.65, 1.68, 3.75, 4.95, "用户可见输入", [
        "自然语言问题：要分析什么、比较什么、解释什么。",
        "引用某条 AI 回复：用于质疑、纠偏、解释证据。",
        "资料范围策略：只看表格、表格+知识库、全上下文。",
        "报告生成范围：当前会话、指定问题、指定周期。"
    ], accent=RUST)
    add_card(slide, 4.75, 1.68, 3.75, 4.95, "系统补充字段", [
        "sessionID / packID / taskID 绑定当前分析对象。",
        "messageID 用于把用户消息和 AI job 串起来。",
        "replyToMessageID 和 quotedMessageSummary 保存引用语境。",
        "contextMode 初始值来自按钮动作，之后可能被路由器降级。"
    ], accent=BLUE)
    add_code(slide, 8.85, 1.68, 3.75, 4.95, [
        "PersistentAIJobPayload(",
        "  prompt: preparedPrompt,",
        "  userMessage: userText,",
        "  messageID: userMessage.id,",
        "  sessionID: session.id,",
        "  packID: pack.id,",
        "  taskID: task?.id,",
        "  contextMode: .fullReanalysis,",
        "  contextSourcePolicy: .tableOnly",
        ")",
    ], title="核心载荷")
    add_footer(slide, page); page += 1

    slide = blank(prs)
    add_title(slide, "2. 会话状态：先落用户消息，再异步执行 AI Job", page=page)
    add_bullets(slide, 0.75, 1.65, 5.45, 4.9, [
        "用户点击发送后，Store 立即 append `AnalysisSessionMessage(role: .user)`，UI 不等待模型。",
        "随后创建 `PersistentAIJob`，写入队列、状态、目标对象和初始 reasoning log。",
        "任务队列负责重试、取消、恢复中断任务、记录 job record。",
        "如果 app 上次退出时 job 仍在 active，启动时恢复为 waiting 并重新调度。",
        "会话状态由 draft/analyzing/waitingForUser/reportReady/archive 统一表达。"
    ], size=13)
    add_card(slide, 6.65, 1.7, 5.65, 1.25, "为什么不同步请求模型", [
        "真实分析会包含表读取、知识库、外部源、Notebook、Harness、修复、报告生成；必须允许取消和恢复。"
    ], accent=RUST)
    add_card(slide, 6.65, 3.25, 5.65, 1.25, "为什么保留 job logs", [
        "logs 是用户看到“正在读取资料 / 正在校验证据 / 自动修正”的来源，也是失败后的复盘记录。"
    ], accent=BLUE)
    add_card(slide, 6.65, 4.8, 5.65, 1.25, "为什么保留 record", [
        "AIJobRecord 是面向分析资产的历史记录：成功、失败、重试次数、错误原因和输出关联。"
    ], accent=SAGE)
    add_footer(slide, page); page += 1

    slide = blank(prs)
    add_title(slide, "3. 路由：决定这轮到底要不要强 Harness", page=page)
    add_table(slide, 0.65, 1.65, 12.0, 4.9,
        ["判断条件", "典型用户问题", "执行路径", "原因"],
        [
            ["轻量任务", "“翻译一下 / 这是什么意思 / 帮我润色”", "quickFollowUp", "不重新读表，避免慢和误拦截。"],
            ["指标关系解释", "“人数和笔数一样，是不是口径问题”", "quickFollowUp", "解释统计口径，不输出未验证数字。"],
            ["表格计算", "“统计 H1 交易金额增长率 / Top5 渠道”", "Harness", "涉及本地数据数字，必须可验证。"],
            ["证据问题", "“结合知识库和外部资料解释变化”", "fullContext", "需要读取知识库/外部参照并附证据。"],
            ["显式重算", "“重新完整分析当前任务”", "fullReanalysis", "用户要求从头读表和重算。"],
        ],
        widths=[2.1, 3.2, 2.0, 4.7],
        font_size=8.7
    )
    add_code(slide, 0.85, 6.15, 11.55, 0.55, [
        "effectiveContextMode(requestedMode, userMessage, hasPreviousAI, cacheMatches) → .cachedFollowUp / .quickFollowUp / .fullReanalysis"
    ])
    add_footer(slide, page); page += 1

    slide = blank(prs)
    add_title(slide, "4. 资料覆盖：AI 在回答前先知道“能看什么”", page=page)
    add_card(slide, 0.65, 1.65, 3.85, 4.95, "Coverage Snapshot", [
        "当前任务选中的表。",
        "表来源：本地 / Tableau。",
        "字段、指标、周期、行数。",
        "外部证据窗口和采集状态。",
        "被跳过或未启用的数据源。"
    ], accent=RUST)
    add_card(slide, 4.8, 1.65, 3.85, 4.95, "用户可见信息", [
        "系统消息说明读取范围。",
        "回答中的“查看证据”。",
        "右侧分析资料面板。",
        "高级审计中可看到 run id、校验结果和 trace。"
    ], accent=BLUE)
    add_card(slide, 8.95, 1.65, 3.35, 4.95, "防止的问题", [
        "模型假装看过没读到的表。",
        "默认拿全周期当某个时间段。",
        "静默忽略未启用的外部源。",
        "报告生成时混入无效会话内容。"
    ], accent=SAGE)
    add_footer(slide, page); page += 1

    slide = blank(prs)
    add_title(slide, "5. 数据接入与结构化：从文件/视图到 Table Manifest", page=page)
    add_bullets(slide, 0.7, 1.56, 5.7, 4.95, [
        "本地导入：CSVParser / ExcelParser / DataImportService 读取 CSV、XLSX、XLS。",
        "Tableau 导入：TableauService 处理登录、视图导出、HTTP 错误解释和 CSV/Crosstab 数据。",
        "结构识别：识别宽表、长表、明细表、Measure Names / Measure Values、指标列、周期列。",
        "字段语义：ReportSemanticInferencer / MetricSemanticExtractionService 推断字段含义、单位、聚合风险。",
        "结果写入 DataPack、ImportedReport、ReportFieldDefinition，后续分析只读结构化对象。"
    ], size=12.1)
    add_code(slide, 6.75, 1.65, 5.55, 4.85, [
        "ImportedReport",
        "  ├─ displayName / sourceType",
        "  ├─ rows / headers / sheet",
        "  ├─ semanticProfile",
        "  ├─ fieldDefinitions",
        "  └─ quality / date range",
        "",
        "TableManifest",
        "  ├─ metricNameColumn",
        "  ├─ metricValueColumn",
        "  ├─ periodColumn",
        "  ├─ aggregationRisk",
        "  └─ normalizedFactTables",
    ], title="内部数据对象")
    add_footer(slide, page); page += 1

    slide = blank(prs)
    add_title(slide, "6. 上下文合成：Prompt 不是拼接，而是分层打包", page=page)
    add_card(slide, 0.65, 1.65, 3.7, 4.95, "业务空间", [
        "国家、时区、币种、业务域。",
        "核心业务链路和指标口径。",
        "分析边界、合规提醒。",
        "默认外部影响因素。"
    ], accent=RUST)
    add_card(slide, 4.55, 1.65, 3.7, 4.95, "会话上下文", [
        "当前目标和用户问题。",
        "上一轮 AI 回答摘要。",
        "引用消息、纠偏记忆。",
        "已生成的报告和证据。"
    ], accent=BLUE)
    add_card(slide, 8.45, 1.65, 3.85, 4.95, "材料上下文", [
        "选中表与字段字典。",
        "Notebook/SQL 计算证据。",
        "知识库、Confluence、Jira、钉钉、本地文件夹。",
        "外部搜索和参照源。"
    ], accent=SAGE)
    add_footer(slide, page); page += 1

    slide = blank(prs)
    add_title(slide, "7. Analysis Harness：强校验路径的核心阶段", page=page)
    steps = [
        ("Manifest", "表结构、周期、指标、聚合风险"),
        ("Intent", "用户到底请求哪些指标/周期/维度"),
        ("Plan", "生成 AnalysisPlan，禁止模型直接算数"),
        ("Execute", "本地执行 sum/count/distinct/growth/rate"),
        ("Report", "模型只解释 verified_results"),
        ("Validate", "数字 trace、引用、边界、缺口"),
    ]
    x = 0.65
    y = 1.74
    for idx, (name, desc) in enumerate(steps):
        row = idx // 3
        col = idx % 3
        add_card(slide, x + col * 4.05, y + row * 2.35, 3.45, 1.7, name, [desc], accent=[RUST, BLUE, SAGE][idx % 3], title_size=15.5, body_size=10.7)
    add_textbox(slide, 0.82, 6.35, 11.2, 0.36, "原则：AI 可以理解意图和解释结论，但关键业务数字必须来自本地 verified_results。", size=15, bold=True, color=INK)
    add_footer(slide, page); page += 1

    slide = blank(prs)
    add_title(slide, "8. 意图解析：先让模型输出 JSON，而不是让它直接回答", page=page)
    add_bullets(slide, 0.72, 1.62, 5.8, 4.92, [
        "AnalysisIntentParser 的系统提示明确：只负责理解用户问题，不负责计算任何业务数字。",
        "输出包含 requested metrics、period、dimensions、filters、derived metrics。",
        "本地校验会阻断未映射指标、缺字段、无法解析周期、AI 意图解析失败。",
        "对于“人数和笔数相同可能是什么原因”这类解释问题，路由层会先降级，避免走强计算链路。",
        "意图解析失败时，不会硬猜本地指标，避免错把“人数”映射到其他人群类指标。"
    ], size=11.7)
    add_code(slide, 6.85, 1.62, 5.45, 4.92, [
        "{",
        '  "requestedMetrics": ["交易人数", "交易笔数"],',
        '  "period": "2026 H1",',
        '  "dimensions": ["渠道"],',
        '  "derivedMetrics": ["人均交易笔数"],',
        '  "filters": [],',
        '  "needsRootCause": false',
        "}",
        "",
        "Validation:",
        "  missingField → fatal",
        "  ambiguousMapping → confirmation",
        "  insufficientData(warning) → warningOnly",
    ], title="意图 JSON 契约")
    add_footer(slide, page); page += 1

    slide = blank(prs)
    add_title(slide, "9. 本地计算：数字从数据结构来，不从自然语言来", page=page)
    add_table(slide, 0.65, 1.58, 12.0, 4.75,
        ["计算类型", "本地处理", "常见风险", "保护机制"],
        [
            ["sum / count", "按 period、metric、dimension 聚合。", "重复记录、粒度不一致。", "duplicateRecordRisk / grainMismatch。"],
            ["distinct count", "识别用户、客户、订单等去重字段。", "把交易笔数当人数。", "distinctCountRisk，必须说明口径。"],
            ["growth / rate", "基期和对比期都由本地结果生成。", "把比例直接求和。", "rateAggregationError。"],
            ["derived metric", "人均、笔均、转化等由 verified result 推导。", "分母缺失或单位冲突。", "formulaMismatch / unit check。"],
            ["root cause", "只有有维度和足够证据时才做贡献拆解。", "无维度时硬归因。", "causalBoundaryRisk。"],
        ],
        widths=[1.9, 3.0, 3.1, 4.0],
        font_size=8.4
    )
    add_textbox(slide, 0.86, 6.45, 11.4, 0.3, "产物：MetricResult[]，包含 rawValue、unit、format、source cells、presentationRole。", size=12.5, bold=True, color=RUST_DARK)
    add_footer(slide, page); page += 1

    slide = blank(prs)
    add_title(slide, "10. SQL / Notebook 证据：把关键计算留下可复盘痕迹", page=page)
    add_card(slide, 0.65, 1.65, 3.7, 4.95, "为什么需要", [
        "业务方会追问“这个数怎么算的”。",
        "AI 报告要能回到表和字段。",
        "跨表联动、派生指标和 period 对齐需要审计路径。",
        "错误时可以定位是解析、计划还是计算问题。"
    ], accent=RUST)
    add_card(slide, 4.55, 1.65, 3.7, 4.95, "记录内容", [
        "SQL/Notebook run id。",
        "输入表、字段、过滤条件。",
        "计算步骤、窗口和 denominator。",
        "输出 MetricResult 与 source cell refs。"
    ], accent=BLUE)
    add_card(slide, 8.45, 1.65, 3.85, 4.95, "展示方式", [
        "回答卡片展示“查看证据”。",
        "右侧证据面板显示 Notebook/SQL 证据。",
        "数字链接可跳转到证据 trace。",
        "导出报告可以复用证据摘要。"
    ], accent=SAGE)
    add_footer(slide, page); page += 1

    slide = blank(prs)
    add_title(slide, "11. 模型调用：流式优先，失败可降级", page=page)
    add_bullets(slide, 0.72, 1.58, 5.75, 4.95, [
        "AIJobQueue 封装重试、validation、correctionPrompt 和 record。",
        "AIStreamingService 请求 chat completions stream，并区分 reasoning progress 与 answer delta。",
        "onProgress 更新“思考过程”；onDelta 更新回答内容。",
        "流式失败时保留已收到的 reasoning，切换普通请求，不让用户卡在半状态。",
        "本地聚合一致性检查失败时，会触发自动重写。"
    ], size=12)
    add_code(slide, 6.75, 1.58, 5.55, 4.95, [
        "runStreamingTextJob(",
        "  prompt, settings, jobType,",
        "  validation: { _ in [] },",
        "  onProgress: { reasoning in",
        "    updateStreamingStatus(reasoning)",
        "  },",
        "  onDelta: { accumulated in",
        "    updateMessageContent(accumulated)",
        "  }",
        ")",
        "",
        "Fallback:",
        "  streaming error → runTextJob(...)",
    ], title="AIJobQueue")
    add_footer(slide, page); page += 1

    slide = blank(prs)
    add_title(slide, "12. 流式 UI：不让 token 直接打爆 SwiftUI", page=page)
    add_card(slide, 0.65, 1.65, 3.85, 4.95, "节流策略", [
        "reasoning 最小间隔：1.3s。",
        "answer 最小间隔：0.9s。",
        "字符增量达到阈值时提前刷新。",
        "savePolicy 使用 deferred，减少磁盘写入。"
    ], accent=RUST)
    add_card(slide, 4.8, 1.65, 3.85, 4.95, "滚动策略", [
        "用户在底部时自动跟随。",
        "手动滚离后暂时暂停 auto-scroll。",
        "分析完成后多次延迟 bottom pin，处理布局延迟。",
        "思考过程文本框内部也跟随底部。"
    ], accent=BLUE)
    add_card(slide, 8.95, 1.65, 3.35, 4.95, "性能保护", [
        "SessionMessageRenderSnapshot。",
        "Markdown / answer parse cache。",
        "Evidence link cache。",
        "长文本阈值短路。"
    ], accent=SAGE)
    add_footer(slide, page); page += 1

    slide = blank(prs)
    add_title(slide, "13. 输出校验：ValidationDecisionEngine 把问题分级", page=page)
    add_table(slide, 0.65, 1.6, 12.0, 4.85,
        ["产品级别", "是否阻断", "典型问题", "处理"],
        [
            ["fatalBlock", "阻断", "未验证数字、缺字段、空结果、公式错误", "不输出未经验证结论，展示需确认资料。"],
            ["needsConfirmation", "需要用户确认", "字段映射歧义、数据契约异常", "提示用户确认表结构或口径。"],
            ["autoRepairable", "不直接阻断", "缺引用、缺方法、证据边界缺失", "自动修复回答结构。"],
            ["warningOnly", "不阻断", "证据不足、外部 claim 风险", "保守输出并显示风险。"],
            ["info", "不阻断", "审计提示", "进入高级审计，不影响主回答。"],
        ],
        widths=[2.0, 1.6, 4.4, 4.0],
        font_size=8.4
    )
    add_textbox(slide, 0.86, 6.43, 11.3, 0.36, "最近的策略调整：insufficientData 的 warning 不再 fatal；只有 error/fatal 才阻断最终输出。", size=12.6, bold=True, color=RUST_DARK)
    add_footer(slide, page); page += 1

    slide = blank(prs)
    add_title(slide, "14. 数字追踪：回答里的数字必须能匹配 verified_results", page=page)
    add_bullets(slide, 0.72, 1.58, 5.75, 4.95, [
        "AnswerNumberTracer 会扫描主回答中的数字、百分比、中文紧凑数值（如 37万）。",
        "matched / approximateMatched 的数字会链接到 MetricResult。",
        "ambiguous / unmatched 的数字进入 validation issues。",
        "引用证据里的外部数字可以被识别为 citation context，不误判成主结论数字。",
        "单位冲突会阻止把“金额”错连到“笔均金额”。"
    ], size=12)
    add_code(slide, 6.75, 1.58, 5.55, 4.95, [
        "AnswerNumberTrace",
        "  rawText: \"37万\"",
        "  status: approximateMatched",
        "  matchedResultID: MetricResult.id",
        "  reason: normalized Chinese compact number",
        "",
        "UI render:",
        "  [37万](nexaflow-evidence://resultID)",
        "",
        "Click → focusMetricResultEvidence(...)",
    ], title="数字到证据的链接")
    add_footer(slide, page); page += 1

    slide = blank(prs)
    add_title(slide, "15. 自动修复：能修的结构问题不打断用户", page=page)
    add_card(slide, 0.65, 1.65, 3.75, 4.95, "修复对象", [
        "直接回答 heading 规范化。",
        "缺 citation 时补入资料证据引用。",
        "过强因果说法降级。",
        "保留限制、读取范围和口径说明。"
    ], accent=RUST)
    add_card(slide, 4.75, 1.65, 3.75, 4.95, "不能修的对象", [
        "缺失本地字段。",
        "指标无法映射。",
        "空结果。",
        "数字与 verified result 不匹配。",
        "不安全 join 或粒度错误。"
    ], accent=BLUE)
    add_card(slide, 8.85, 1.65, 3.75, 4.95, "用户看到什么", [
        "严重问题：需要确认分析资料。",
        "轻风险：回答保守输出。",
        "审计提示：收进分析资料面板。",
        "不会展示内部 raw validation code 给业务用户。"
    ], accent=SAGE)
    add_footer(slide, page); page += 1

    slide = blank(prs)
    add_title(slide, "16. 回答组织：直接回答和依据分离", page=page)
    add_bullets(slide, 0.72, 1.58, 5.75, 4.95, [
        "AnalysisAnswerPresentation 解析 `## 直接回答你的问题` 等标题。",
        "聊天卡片优先显示直接回答，资料、计算证据、限制说明收进分析资料。",
        "复制按钮默认复制直接回答，菜单里可复制完整原文。",
        "没有结构化标题的老消息仍按 legacy Markdown 渲染。",
        "这样避免聊天区被证据和审计文本撑爆。"
    ], size=12)
    add_code(slide, 6.75, 1.58, 5.55, 4.95, [
        "## 直接回答你的问题",
        "人数和笔数相同，通常意味着...",
        "",
        "## 本地事实",
        "- 已读取表 A / 表 B",
        "",
        "## 计算证据",
        "- verified result refs",
        "",
        "## 限制与不确定性",
        "- 未覆盖外部来源...",
    ], title="回答 Markdown 契约")
    add_footer(slide, page); page += 1

    slide = blank(prs)
    add_title(slide, "17. Markdown 渲染：表格、长文本和证据链接都有专门处理", page=page)
    add_card(slide, 0.65, 1.65, 3.75, 4.95, "MarkdownMessageRenderer", [
        "解析 text/table block。",
        "缓存 block 和 AttributedString。",
        "长文本进入 LongTextPreview。",
        "表格只渲染前 60 行。"
    ], accent=RUST)
    add_card(slide, 4.75, 1.65, 3.75, 4.95, "表格优化", [
        "列宽在 parse 阶段计算一次。",
        "水平滚动避免撑坏页面。",
        "截断提示告诉用户完整内容仍在原文。",
        "交替行背景提高可读性。"
    ], accent=BLUE)
    add_card(slide, 8.85, 1.65, 3.75, 4.95, "证据链接优化", [
        "只在非流式、展开状态插入链接。",
        "按 message/run/traces/content 指纹缓存。",
        "点击链接打开右侧证据。",
        "避免每次 body 重扫全文。"
    ], accent=SAGE)
    add_footer(slide, page); page += 1

    slide = blank(prs)
    add_title(slide, "18. 证据与审计面板：回答不是终点，是可复盘资产", page=page)
    add_bullets(slide, 0.72, 1.56, 5.75, 5.0, [
        "每条 AI 回答可以携带 AnalysisSessionEvidence。",
        "Evidence 类型包括数据覆盖、计算证据、Analysis Harness 审计、外部资料、知识库等。",
        "右侧分析资料面板聚合展示读取范围、表字段、Notebook、Metric Semantic、证据 trace。",
        "AIJobRecord 保留执行状态、日志、错误和目标对象。",
        "纠偏和采纳动作会把回答转成 correction memory 或知识库资产。"
    ], size=12)
    add_table(slide, 6.65, 1.62, 5.75, 4.95,
        ["对象", "用途"],
        [
            ["AnalysisSessionEvidence", "回答级证据容器。"],
            ["AnalysisHarnessRun", "强校验 run 的完整审计。"],
            ["AnalysisCoverageSnapshot", "本轮读取范围和限制。"],
            ["AnalysisNotebookRun", "计算过程证据。"],
            ["CorrectionMemory", "后续分析复用的纠偏规则。"],
        ],
        widths=[2.4, 3.35],
        font_size=9.0
    )
    add_footer(slide, page); page += 1

    slide = blank(prs)
    add_title(slide, "19. 持久化与密钥：用户自己的 API，配置随 workspace 保存", page=page)
    add_card(slide, 0.65, 1.65, 3.75, 4.95, "workspace.json", [
        "ProductWorkspace 统一保存业务空间、数据包、会话、AI 设置。",
        "AISettings 包含 endpoint、model、apiKey、systemPrompt。",
        "按用户偏好不使用 Keychain。"
    ], accent=RUST)
    add_card(slide, 4.75, 1.65, 3.75, 4.95, "安全边界", [
        "请求时只放到 Authorization header。",
        "日志和 debug snapshot 不输出 raw key。",
        "错误信息避免回显密钥。",
        "外部源 key 也走独立设置。"
    ], accent=BLUE)
    add_card(slide, 8.85, 1.65, 3.75, 4.95, "可靠性", [
        "workspace 读写有 corrupt backup。",
        "API key 保存/读取有回归测试。",
        "任务取消错误不会误重试。",
        "job logs 做上限裁剪。"
    ], accent=SAGE)
    add_footer(slide, page); page += 1

    slide = blank(prs)
    add_title(slide, "20. 错误、降级和重试：失败也要有可解释状态", page=page)
    add_table(slide, 0.65, 1.6, 12.0, 5.0,
        ["失败点", "处理", "用户看到的状态"],
        [
            ["流式连接失败", "保留 reasoning，降级 runTextJob。", "流式输出已降级为普通请求。"],
            ["Harness 基础设施异常", "移除临时 assistant 消息，回退旧 AI 直答链路。", "本地校验回退。"],
            ["AI 意图解析失败", "阻断强计算，不本地硬猜。", "需要确认分析资料。"],
            ["Tableau 502 / HTML 错误页", "转成可读错误和技术详情。", "Tableau 服务暂时不可用。"],
            ["用户取消", "停止 job，恢复会话可操作状态。", "已取消当前分析。"],
        ],
        widths=[2.5, 5.1, 4.4],
        font_size=8.8
    )
    add_footer(slide, page); page += 1

    slide = blank(prs)
    add_title(slide, "21. 性能处理：让长回答和长证据不会拖垮 UI", page=page)
    add_bullets(slide, 0.72, 1.55, 5.75, 5.0, [
        "SessionMessageRenderSnapshot 让 Equatable 只比较 UI 必要字段和有界文本指纹。",
        "latestMessageRenderRevision 不再对长字符串 hash，使用 id、长度、状态、证据数。",
        "Markdown block、AttributedString、AnalysisAnswerPresentation 都有 NSCache。",
        "证据数字链接有缓存，避免反复替换大文本。",
        "PlainScrollableTextView 自适应高度，避免固定大空白和内部滚动错位。"
    ], size=12)
    add_code(slide, 6.75, 1.55, 5.55, 5.0, [
        "NEXAFLOW_PERF_TRACE=1",
        "",
        "markdown.blocks chars=2041 took 3.2ms",
        "markdown.attributed chars=812 took 1.1ms",
        "analysisAnswer.parse chars=2231 took 0.7ms",
        "message.linkedEvidenceContent traces=9 took 2.8ms",
        "",
        "Default: trace disabled, no user text logged.",
    ], title="开发期性能 trace")
    add_footer(slide, page); page += 1

    slide = blank(prs)
    add_title(slide, "22. 测试矩阵：每次改链路都要覆盖四层", page=page)
    add_table(slide, 0.65, 1.58, 12.0, 5.0,
        ["层级", "命令 / 场景", "验证目标"],
        [
            ["编译", "swift build --disable-sandbox --product IterationPilot", "所有 Swift 类型和 app target 可构建。"],
            ["非 GUI 回归", "swift run --disable-sandbox IterationPilotRegressionTests", "Harness、路由、校验、解析、可靠性。"],
            ["打包级回归", "./script/non_gui_regression.sh", "debug/release/universal binary。"],
            ["App verify", "./script/build_and_run.sh --verify", "窗口可见、debug snapshots OK。"],
            ["Live smoke", "NEXAFLOW_LIVE_AI_SMOKE=1 ...", "真实流式 AI 和 Harness 链路。"],
            ["手动验收", "长思考过程、完成后滚动、解释类追问", "用户看到的核心体验。"],
        ],
        widths=[2.0, 5.2, 4.8],
        font_size=8.6
    )
    add_footer(slide, page); page += 1

    slide = blank(prs)
    add_title(slide, "23. 代码地图：主要文件和职责", page=page)
    add_table(slide, 0.65, 1.48, 12.0, 5.28,
        ["模块", "关键文件", "职责"],
        [
            ["UI 会话", "AnalysisSessionsView / AnalysisSessionChatView", "消息列表、滚动、消息卡片、证据链接、报告入口。"],
            ["状态编排", "AnalysisSessionStoreActions / AIJobStoreActions", "发送消息、建 job、调度、状态更新、缓存上下文。"],
            ["AI 请求", "AIJobQueue / AIStreamingService / AIServiceResponseParser", "流式请求、普通请求、重试、解析。"],
            ["Harness", "AnalysisHarnessServices / ValidationDecisionEngine", "意图解析、计划、本地计算、报告、校验、修复。"],
            ["数据事实层", "CSVParser / ExcelParser / TableauService / TableContextPackageBuilder", "导入、结构化、表上下文打包。"],
            ["证据和展示", "AnalysisAnswerPresentation / MarkdownMessageRenderer / DebugSnapshotExporter", "回答拆分、Markdown 渲染、debug snapshot。"],
        ],
        widths=[1.8, 4.8, 5.4],
        font_size=8.2
    )
    add_footer(slide, page); page += 1

    slide = blank(prs)
    add_title(slide, "24. 关键设计取舍和下一步", page=page)
    add_card(slide, 0.65, 1.65, 3.75, 4.95, "已做取舍", [
        "用户自己的 API key 保存到 workspace。",
        "解释类问题优先可用性，不被强校验误拦截。",
        "数字输出保持强校验。",
        "UI 优先做缓存和节流，避免大重构打断功能。"
    ], accent=RUST)
    add_card(slide, 4.75, 1.65, 3.75, 4.95, "仍需优化", [
        "拆分 AnalysisHarnessServices。",
        "拆分 AnalysisSessionsView。",
        "PersistentAIJob fingerprint 避免大 payload 编码。",
        "更完整的可视化 trace timeline。"
    ], accent=BLUE)
    add_card(slide, 8.85, 1.65, 3.75, 4.95, "对外表达", [
        "NexaFlow 不是套壳大模型。",
        "它是可审计的业务分析工作台。",
        "核心价值在于事实层、证据层、校验层和报告资产沉淀。",
        "AI 负责解释，本地系统负责约束。"
    ], accent=SAGE)
    add_footer(slide, page)

    return prs


def main():
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs = build_deck()
    prs.save(OUT)
    LATEST.write_text(str(OUT) + "\n", encoding="utf-8")
    print(OUT)


if __name__ == "__main__":
    main()
